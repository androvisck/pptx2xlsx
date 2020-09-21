"""
contar o número de slides
ler OS, SW, name, version, etc
listar as info em cada linha do csv
listar cada slide de issue após as info em cada linha
"""

from pptx import Presentation

local_pptxFileList = ["test.pptx"]
text_runs = []
count=0
for i in local_pptxFileList:
    ppt = Presentation(i)
    for slide in ppt.slides:
        count=count+1
        for shape in slide.shapes:  # para cada caixa de texto inserida num slide
            if shape.has_text_frame:  # verifica se tem uma caixa de texto
                text_runs.append(shape.text)  # adiciona o conteúdo das caixas de texto numa lista
                print(shape.text)
print('Número de slides= ',count)

print("\n")
print("Imprimir a Lista: ",text_runs)

print("\nSeparando os Campos:")
titulo=text_runs[0]
print(titulo)
y=text_runs[1]
x = y.split('\n')
#print(x)
nome=x[0]
os=x[1]
versao=x[2]
#print(nome,os,versao)
nome = nome.split('Nome:')
os = os.split('OS:')
versao = versao.split('Versão:')
nome=nome[1]
os=os[1]
versao=versao[1]

nome=" ".join(nome.split())
os=" ".join(os.split())
versao=" ".join(versao.split())
print(nome)
print(os)
print(versao)
print('\n')

# importando pandas como pd
import pandas as pd
head=['titulo','nome','sw','versao']
# Construindo o DataFrame a partir de uma lista
df = pd.DataFrame(head)
pd.Index(head)
print(df)

import xlsxwriter

new_list = [['titulo','nome','sw','versao'], [1, 2, 3, 4, 5, 6]]

with xlsxwriter.Workbook('test.xlsx') as workbook:
    worksheet = workbook.add_worksheet()

    for row_num, data in enumerate(new_list):
        worksheet.write_row(row_num, 0, data)